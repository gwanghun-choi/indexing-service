import os
import asyncio
import time
from pathlib import Path
from typing import List, Dict
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import logging

from app.parser.base import ParserInterface

logger = logging.getLogger(__name__)


class PptParser(ParserInterface):
    """PowerPoint 파일 파서 클래스"""
    
    def __init__(self):
        """PptParser 초기화"""
        self.supported_extensions = [".ppt", ".pptx"]

    async def parsing(
        self, file_path: str, filename: str = None
    ) -> List[Dict]:
        """PPT/PPTX 파일을 파싱하여 마크다운 텍스트와 이미지를 추출"""
        try:
            logger.info(f"PPT 파일 파싱 시작: file_path={file_path}, filename={filename}")

            # 파일 확장자 확인 - filename에서 먼저 확인
            if filename:
                file_ext = Path(filename).suffix.lower()
            else:
                file_ext = Path(file_path).suffix.lower()

            if not file_ext:
                raise ValueError("파일 확장자를 확인할 수 없습니다")

            if file_ext not in self.supported_extensions:
                raise ValueError(f"지원되지 않는 파일 형식: {file_ext}")

            result = await self._read_pptx(file_path)

            logger.info(f"PPT 파일 파싱 완료: {len(result)}개 슬라이드 처리됨")
            return result

        except Exception as e:
            logger.error(f"PPT 파싱 중 오류 발생: {e}")
            raise

    async def _read_pptx(self, file_path: str) -> List[Dict]:
        """PowerPoint 파일 비동기 파싱"""
        return await asyncio.to_thread(self._pptx_sync, file_path)

    @staticmethod
    def _pptx_sync(file_path: str) -> List[Dict]:
        """PowerPoint 파일 동기 파싱"""
        try:
            prs = Presentation(file_path)
            result_list: List[Dict] = []

            # 임시 이미지 폴더 생성
            project_root = Path(__file__).parent.parent.parent
            temp_images_dir = project_root / "uploads" / "temp_images"
            temp_images_dir.mkdir(parents=True, exist_ok=True)

            ppt_name = Path(file_path).stem
            timestamp = str(int(time.time() * 1000))
            temp_dir = temp_images_dir / f"{ppt_name}_{timestamp}"
            temp_dir.mkdir(exist_ok=True)
            temp_dir = str(temp_dir)

            # 각 슬라이드 처리
            for slide_idx, slide in enumerate(prs.slides):
                slide_number = slide_idx + 1
                markdown_content = []
                images_info = []
                img_counter = 0

                # 슬라이드의 모든 shape 처리
                for shape in slide.shapes:
                    # 텍스트 처리
                    if shape.has_text_frame:
                        text = PptParser._extract_text_from_shape(shape)
                        if text.strip():
                            # 제목 처리
                            if shape == slide.shapes.title:
                                markdown_content.append(
                                    f"# 슬라이드 {slide_number}: {text}"
                                )
                            else:
                                markdown_content.append(text)

                    # 테이블 처리
                    elif shape.has_table:
                        table_md = PptParser._table_to_markdown(shape.table)
                        if table_md.strip():
                            markdown_content.append(table_md)

                    # 차트 처리
                    elif shape.has_chart:
                        chart_title = "차트"
                        if (
                            shape.chart.has_title
                            and shape.chart.chart_title.has_text_frame
                        ):
                            chart_title = shape.chart.chart_title.text_frame.text
                        markdown_content.append(f"[차트: {chart_title}]")

                    # 이미지 처리
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_counter += 1
                        try:
                            # 이미지 추출 및 저장
                            image = shape.image
                            image_bytes = image.blob

                            # 이미지 파일 확장자 결정
                            ext = image.ext
                            if not ext:
                                ext = "png"

                            img_filename = (
                                f"slide_{slide_number}_img_{img_counter}.{ext}"
                            )
                            img_path = os.path.join(temp_dir, img_filename)

                            with open(img_path, "wb") as img_file:
                                img_file.write(image_bytes)

                            images_info.append(
                                {
                                    "filename": img_filename,
                                    "path": img_path,
                                    "width": shape.width,
                                    "height": shape.height,
                                    "type": "embedded_image",
                                    "page": slide_number,
                                    "note": f"슬라이드 {slide_number}의 이미지 {img_counter}",
                                }
                            )

                        except Exception as e:
                            logger.warning(f"이미지 추출 중 오류: {e}")
                            continue

                # 슬라이드 텍스트 결합
                slide_text = (
                    "\n\n".join(markdown_content)
                    if markdown_content
                    else f"# 슬라이드 {slide_number}\n\n*(빈 슬라이드)*"
                )

                result_list.append(
                    {
                        "page_number": slide_number,
                        "text": slide_text,
                        "images": images_info,
                        "temp_dir": temp_dir,
                    }
                )

            return result_list

        except Exception as e:
            logger.error(f"PPT 파싱 중 오류 발생: {e}")
            raise

    @staticmethod
    def _extract_text_from_shape(shape):
        """Shape에서 텍스트 추출 및 포맷팅"""
        text_runs = []

        for paragraph in shape.text_frame.paragraphs:
            para_text = []

            for run in paragraph.runs:
                text = run.text

                # 볼드와 이탤릭 처리
                if run.font.bold and run.font.italic:
                    text = f"***{text}***"
                elif run.font.bold:
                    text = f"**{text}**"
                elif run.font.italic:
                    text = f"*{text}*"

                para_text.append(text)

            full_para_text = "".join(para_text).strip()

            if full_para_text:
                # 불릿 포인트 처리
                if paragraph.level > 0:
                    indent = "  " * (paragraph.level - 1)
                    text_runs.append(f"{indent}- {full_para_text}")
                else:
                    text_runs.append(full_para_text)

        return "\n".join(text_runs)

    @staticmethod
    def _table_to_markdown(table):
        """테이블을 마크다운 형식으로 변환"""
        if not table.rows:
            return ""

        markdown_lines = []

        # 헤더 행
        header_cells = []
        for cell in table.rows[0].cells:
            cell_text = cell.text.strip().replace("\n", "<br>")
            header_cells.append(cell_text)

        markdown_lines.append("| " + " | ".join(header_cells) + " |")
        markdown_lines.append("| " + " | ".join(["---"] * len(header_cells)) + " |")

        # 데이터 행
        for i in range(1, len(table.rows)):
            row = table.rows[i]
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", "<br>")
                row_cells.append(cell_text)
            markdown_lines.append("| " + " | ".join(row_cells) + " |")

        return "\n".join(markdown_lines)
